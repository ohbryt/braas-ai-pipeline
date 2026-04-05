"""
Report Generator
================

Generate experimental reports in multiple formats: PDF, LaTeX, JSON, 
PowerPoint, and Markdown.
"""

from __future__ import annotations

import json
import os
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any

from braas.core.enums import ExperimentType
from braas.core.exceptions import AnalysisError
from braas.core.models import ExperimentResult


# -----------------------------------------------------------------------------
# Data Classes
# -----------------------------------------------------------------------------

@dataclass
class ReportMetadata:
    """Metadata for a generated report."""
    title: str
    author: str = "BRaaS Pipeline"
    institution: str = ""
    version: str = "1.0"
    generated_at: datetime = field(default_factory=datetime.now)


# -----------------------------------------------------------------------------
# Report Generator
# -----------------------------------------------------------------------------

class ReportGenerator:
    """Generate reports in multiple formats from experimental results.
    
    Supported formats:
    - PDF: Using reportlab
    - LaTeX: Direct .tex file generation
    - JSON: Structured data export
    - PowerPoint: Using python-pptx
    - Markdown: Human-readable text format
    """
    
    def __init__(self, output_dir: Path | None = None) -> None:
        """Initialize the report generator.
        
        Args:
            output_dir: Directory to save reports. Defaults to outputs/reports/
        """
        self._output_dir = output_dir or self._get_default_output_dir()
        self._output_dir.mkdir(parents=True, exist_ok=True)
    
    def _get_default_output_dir(self) -> Path:
        """Get default output directory."""
        return Path.home() / "braas-ai-pipeline" / "outputs" / "reports"
    
    def _generate_filename(self, prefix: str, extension: str) -> str:
        """Generate unique filename."""
        import uuid
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        unique_id = uuid.uuid4().hex[:6]
        return f"{prefix}_{timestamp}_{unique_id}.{extension}"
    
    def generate_pdf_report(self, experiment_result: ExperimentResult) -> str:
        """Generate PDF report from experiment results.
        
        Sections:
        - Cover page
        - Executive Summary
        - Methods
        - Results (with figures)
        - Statistical Analysis
        - AI Insights
        - References
        
        Args:
            experiment_result: ExperimentResult object with analysis data
        
        Returns:
            Path to saved PDF file
        """
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import (
                SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
                Image, PageBreak, ListFlowable, ListItem
            )
            from reportlab.lib import colors
        except ImportError:
            raise AnalysisError(
                message="reportlab is required for PDF generation. Install with: pip install reportlab",
                analysis_type="pdf_generation"
            )
        
        filepath = self._output_dir / self._generate_filename("report", "pdf")
        
        # Create document
        doc = SimpleDocTemplate(
            str(filepath),
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )
        
        # Styles
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            alignment=1  # Center
        )
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=14,
            spaceBefore=20,
            spaceAfter=10,
            textColor=colors.darkblue
        )
        body_style = styles['BodyText']
        
        # Build content
        story = []
        
        # Cover Page
        story.append(Spacer(1, 2 * inch))
        story.append(Paragraph("BRaaS Pipeline Report", title_style))
        story.append(Spacer(1, 0.5 * inch))
        story.append(Paragraph(f"Experiment ID: {experiment_result.experiment_id}", body_style))
        story.append(Paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}", body_style))
        story.append(Paragraph(f"Quality Score: {experiment_result.quality_score:.2%}", body_style))
        story.append(PageBreak())
        
        # Executive Summary
        story.append(Paragraph("Executive Summary", heading_style))
        summary_text = self._get_summary_text(experiment_result)
        story.append(Paragraph(summary_text, body_style))
        story.append(Spacer(1, 0.3 * inch))
        
        # Methods
        story.append(Paragraph("Methods", heading_style))
        story.append(Paragraph(
            "The experiment was conducted following standard operating procedures. "
            "Data was collected using automated instrumentation and processed using "
            "the BRaaS AI pipeline analysis engine.",
            body_style
        ))
        story.append(Spacer(1, 0.3 * inch))
        
        # Results
        story.append(Paragraph("Results", heading_style))
        
        if experiment_result.summary:
            for key, value in experiment_result.summary.items():
                story.append(Paragraph(f"<b>{key}:</b> {value}", body_style))
        
        story.append(Spacer(1, 0.3 * inch))
        
        # Statistical Analysis
        if experiment_result.statistical_tests:
            story.append(Paragraph("Statistical Analysis", heading_style))
            for test in experiment_result.statistical_tests:
                test_name = test.get('test_name', 'Test')
                p_value = test.get('p_value', 'N/A')
                story.append(Paragraph(f"{test_name}: p = {p_value}", body_style))
        
        story.append(Spacer(1, 0.3 * inch))
        
        # Figures
        if experiment_result.plots_uris:
            story.append(Paragraph("Figures", heading_style))
            for plot_uri in experiment_result.plots_uris:
                if os.path.exists(plot_uri):
                    try:
                        img = Image(plot_uri, width=4*inch, height=3*inch)
                        story.append(img)
                        story.append(Spacer(1, 0.2 * inch))
                    except Exception:
                        story.append(Paragraph(f"[Figure: {Path(plot_uri).name}]", body_style))
        
        # AI Insights
        story.append(Paragraph("AI Insights", heading_style))
        story.append(Paragraph(
            "The BRaaS AI analysis pipeline identified key patterns in the data. "
            "Further investigation is recommended for any significant findings.",
            body_style
        ))
        
        # References
        story.append(Paragraph("References", heading_style))
        story.append(Paragraph(
            "1. BRaaS Pipeline Documentation<br/>"
            "2. Experimental Protocol Records<br/>"
            "3. Instrument Calibration Logs",
            body_style
        ))
        
        # Build PDF
        doc.build(story)
        
        return str(filepath)
    
    def _get_summary_text(self, experiment_result: ExperimentResult) -> str:
        """Generate summary text from experiment result."""
        if experiment_result.passed_qc:
            quality_status = "PASSED"
        else:
            quality_status = "DID NOT PASS"
        
        return (
            f"This report summarizes the analysis of experiment {experiment_result.experiment_id}. "
            f"The experiment {quality_status} quality control thresholds with a quality score of "
            f"{experiment_result.quality_score:.2%}. "
            f"Key metrics and statistical analyses are detailed in the following sections."
        )
    
    def generate_latex_report(self, experiment_result: ExperimentResult) -> str:
        """Generate LaTeX report from experiment results.
        
        Args:
            experiment_result: ExperimentResult object
        
        Returns:
            Path to saved .tex file
        """
        filepath = self._output_dir / self._generate_filename("report", "tex")
        
        content = self._generate_latex_content(experiment_result)
        
        with open(filepath, 'w') as f:
            f.write(content)
        
        return str(filepath)
    
    def _generate_latex_content(self, experiment_result: ExperimentResult) -> str:
        """Generate LaTeX document content."""
        sections = []
        
        # Document header
        sections.append(r"\documentclass[12pt]{article}")
        sections.append(r"\usepackage[utf8]{inputenc}")
        sections.append(r"\usepackage{graphicx}")
        sections.append(r"\usepackage{hyperref}")
        sections.append(r"\usepackage{amsmath}")
        sections.append(r"\usepackage{booktabs}")
        sections.append(r"\usepackage[left=1in,right=1in,top=1in,bottom=1in]{geometry}")
        sections.append(r"\usepackage{times}")
        sections.append(r"\setlength{\parindent}{0.5in}")
        sections.append(r"\setlength{\parskip}{0.5em}")
        sections.append(r"\begin{document}")
        
        # Title
        sections.append(r"\begin{titlepage}")
        sections.append(r"\centering")
        sections.append(r"\vspace*{2cm}")
        sections.append(r"{\LARGE\bfseries BRaaS Pipeline Report\par}")
        sections.append(r"\vspace{1.5cm}")
        sections.append(f"{{\\large Experiment ID: {experiment_result.experiment_id}\\par}}")
        sections.append(f"{{\\large Date: {datetime.now().strftime('%Y-%m-%d')}\\par}}")
        sections.append(r"\vspace{1cm}")
        sections.append(f"Quality Score: {experiment_result.quality_score:.2%}\\par")
        sections.append(r"\end{titlepage}")
        sections.append(r"\newpage")
        
        # Executive Summary
        sections.append(r"\section{Executive Summary}")
        sections.append(self._get_summary_text(experiment_result).replace('. ', r'. \par' + '\n'))
        
        # Methods
        sections.append(r"\section{Methods}")
        sections.append(r"Data was collected and processed following standard protocols.")
        
        # Results
        sections.append(r"\section{Results}")
        if experiment_result.summary:
            sections.append(r"\begin{itemize}")
            for key, value in experiment_result.summary.items():
                sections.append(f"\\item {{{key}: {value}}}")
            sections.append(r"\end{itemize}")
        
        # Statistical Tests
        if experiment_result.statistical_tests:
            sections.append(r"\section{Statistical Analysis}")
            sections.append(r"\begin{table}[h]")
            sections.append(r"\centering")
            sections.append(r"\begin{tabular}{ll}")
            sections.append(r"\toprule")
            sections.append(r"Test & p-value \\")
            sections.append(r"\midrule")
            for test in experiment_result.statistical_tests:
                test_name = test.get('test_name', 'Test')
                p_value = test.get('p_value', 'N/A')
                sections.append(f"{test_name} & {p_value} \\\\")
            sections.append(r"\bottomrule")
            sections.append(r"\end{tabular}")
            sections.append(r"\end{table}")
        
        # Figures
        if experiment_result.plots_uris:
            sections.append(r"\section{Figures}")
            for i, plot_uri in enumerate(experiment_result.plots_uris):
                sections.append(f"\\begin{{figure}}[h]")
                sections.append(f"\\centering")
                sections.append(f"\\includegraphics[width=0.8\\textwidth]{{{plot_uri}}}")
                sections.append(f"\\caption{{Figure {i+1}}}")
                sections.append(f"\\end{{figure}}")
        
        # References
        sections.append(r"\section{References}")
        sections.append(r"\begin{enumerate}")
        sections.append(r"\item BRaaS Pipeline Documentation")
        sections.append(r"\item Experimental Protocol Records")
        sections.append(r"\end{enumerate}")
        
        sections.append(r"\end{document}")
        
        return '\n'.join(sections)
    
    def generate_json_report(self, experiment_result: ExperimentResult) -> dict[str, Any]:
        """Generate JSON report from experiment results.
        
        Args:
            experiment_result: ExperimentResult object
        
        Returns:
            Dictionary representation of the report
        """
        report = {
            "report_metadata": {
                "generated_at": datetime.now().isoformat(),
                "version": "1.0",
                "generator": "BRaaS Pipeline"
            },
            "experiment_info": {
                "experiment_id": experiment_result.experiment_id,
                "quality_score": experiment_result.quality_score,
                "passed_qc": experiment_result.passed_qc
            },
            "summary": experiment_result.summary,
            "statistical_tests": experiment_result.statistical_tests,
            "plots_uris": experiment_result.plots_uris,
            "ml_predictions": experiment_result.ml_predictions,
            "notes": experiment_result.notes
        }
        
        return report
    
    def generate_slide_deck(self, experiment_result: ExperimentResult) -> str:
        """Generate PowerPoint slide deck from experiment results.
        
        Args:
            experiment_result: ExperimentResult object
        
        Returns:
            Path to saved .pptx file
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
        except ImportError:
            raise AnalysisError(
                message="python-pptx is required for PowerPoint generation. Install with: pip install python-pptx",
                analysis_type="pptx_generation"
            )
        
        filepath = self._output_dir / self._generate_filename("presentation", "pptx")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "BRaaS Pipeline Report"
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = f"Experiment: {experiment_result.experiment_id}"
        sub_para.font.size = Pt(24)
        sub_para.alignment = PP_ALIGN.CENTER
        
        # Executive Summary slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Executive Summary"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        content_frame = content_box.text_frame
        content_para = content_frame.paragraphs[0]
        content_para.text = f"Quality Score: {experiment_result.quality_score:.2%}"
        content_para.font.size = Pt(20)
        
        if experiment_result.summary:
            for key, value in experiment_result.summary.items():
                p = content_frame.add_paragraph()
                p.text = f"{key}: {value}"
                p.font.size = Pt(18)
        
        # Results slide
        if experiment_result.summary:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = "Key Results"
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
            content_frame = content_box.text_frame
            for i, (key, value) in enumerate(experiment_result.summary.items()):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                p.text = f"• {key}: {value}"
                p.font.size = Pt(20)
        
        # Figures slide
        if experiment_result.plots_uris:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = "Results Figures"
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            
            for i, plot_uri in enumerate(experiment_result.plots_uris[:4]):  # Limit to 4 figures
                if os.path.exists(plot_uri):
                    try:
                        left = Inches(0.5 + (i % 2) * 4.5)
                        top = Inches(1.5 + (i // 2) * 2.8)
                        slide.shapes.add_picture(
                            plot_uri,
                            left, top,
                            width=Inches(4)
                        )
                    except Exception:
                        pass
        
        # Conclusions slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Conclusions"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        content_frame = content_box.text_frame
        p = content_frame.paragraphs[0]
        p.text = "• Analysis completed successfully"
        p.font.size = Pt(20)
        
        if experiment_result.passed_qc:
            p = content_frame.add_paragraph()
            p.text = "• Results meet quality control criteria"
            p.font.size = Pt(20)
        
        p = content_frame.add_paragraph()
        p.text = "• See detailed report for complete analysis"
        p.font.size = Pt(20)
        
        prs.save(str(filepath))
        
        return str(filepath)
    
    def generate_markdown_report(self, experiment_result: ExperimentResult) -> str:
        """Generate Markdown report from experiment results.
        
        Args:
            experiment_result: ExperimentResult object
        
        Returns:
            Markdown-formatted string
        """
        lines = []
        
        # Title
        lines.append("# BRaaS Pipeline Report")
        lines.append("")
        lines.append(f"**Experiment ID:** {experiment_result.experiment_id}")
        lines.append(f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        lines.append(f"**Quality Score:** {experiment_result.quality_score:.2%}")
        lines.append("")
        lines.append("---")
        lines.append("")
        
        # Executive Summary
        lines.append("## Executive Summary")
        lines.append("")
        lines.append(self._get_summary_text(experiment_result))
        lines.append("")
        
        # Methods
        lines.append("## Methods")
        lines.append("")
        lines.append("The experiment was conducted following standard operating procedures. "
                    "Data was collected using automated instrumentation and processed using "
                    "the BRaaS AI pipeline analysis engine.")
        lines.append("")
        
        # Results
        lines.append("## Results")
        lines.append("")
        
        if experiment_result.summary:
            for key, value in experiment_result.summary.items():
                lines.append(f"**{key}:** {value}")
            lines.append("")
        
        # Statistical Analysis
        if experiment_result.statistical_tests:
            lines.append("## Statistical Analysis")
            lines.append("")
            lines.append("| Test | p-value |")
            lines.append("|------|---------|")
            for test in experiment_result.statistical_tests:
                test_name = test.get('test_name', 'Test')
                p_value = test.get('p_value', 'N/A')
                lines.append(f"| {test_name} | {p_value} |")
            lines.append("")
        
        # Figures
        if experiment_result.plots_uris:
            lines.append("## Figures")
            lines.append("")
            for i, plot_uri in enumerate(experiment_result.plots_uris):
                lines.append(f"### Figure {i+1}")
                lines.append(f"![Result Figure]({plot_uri})")
                lines.append("")
        
        # AI Insights
        lines.append("## AI Insights")
        lines.append("")
        lines.append("The BRaaS AI analysis pipeline identified key patterns in the data. "
                    "Further investigation is recommended for any significant findings.")
        lines.append("")
        
        # References
        lines.append("## References")
        lines.append("")
        lines.append("1. BRaaS Pipeline Documentation")
        lines.append("2. Experimental Protocol Records")
        lines.append("3. Instrument Calibration Logs")
        lines.append("")
        
        return "\n".join(lines)
